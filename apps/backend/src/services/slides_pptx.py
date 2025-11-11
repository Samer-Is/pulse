"""Slides generation service using python-pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt
import io

from ..schemas.slides import SlidesGenerateRequest, SlideContent
from ..core.logging import get_logger

logger = get_logger(__name__)


def generate_slides_pptx(request: SlidesGenerateRequest) -> io.BytesIO:
    """
    Generate presentation slides as PPTX file.
    
    Args:
        request: Slides generation request with all data
    
    Returns:
        BytesIO buffer containing PPTX file
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = request.presentation_title
    subtitle.text = f"By {request.author}"
    
    # Content Slides
    for slide_content in sorted(request.slides, key=lambda x: x.order):
        # Use different layouts based on layout type
        if slide_content.layout == "title-only":
            slide_layout = prs.slide_layouts[5]  # Blank
        elif slide_content.layout == "two-column":
            slide_layout = prs.slide_layouts[3]  # Two content
        else:
            slide_layout = prs.slide_layouts[1]  # Title and content
        
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        if slide.shapes.title:
            slide.shapes.title.text = slide_content.title
        
        # Add content
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.text = slide_content.content
            
            # Format text
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(18)
    
    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    logger.info(f"Slides generated: title={request.presentation_title}, slides={len(request.slides)}")
    
    return buffer

