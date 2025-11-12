"""Slide Generator Service - Creates PPTX and PDF presentations."""

import io
from typing import List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from playwright.async_api import async_playwright

from ..schemas.slides import SlideContent


class SlideGenerator:
    """Generate presentations in PPTX and PDF formats."""

    def __init__(self):
        """Initialize slide generator."""
        self.title_color = RGBColor(31, 78, 121)  # Professional blue
        self.text_color = RGBColor(51, 51, 51)  # Dark gray

    async def generate_pptx(
        self,
        title: str,
        slides: List[SlideContent],
        template: str = "modern"
    ) -> bytes:
        """
        Generate presentation in PPTX format.
        
        Args:
            title: Presentation title
            slides: List of slide contents
            template: Template style
            
        Returns:
            PPTX file as bytes
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add title slide
        self._add_title_slide(prs, title)
        
        # Add content slides
        for slide_content in slides:
            self._add_content_slide(prs, slide_content)
        
        # Save to bytes
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        return file_stream.getvalue()

    def _add_title_slide(self, prs: Presentation, title: str) -> None:
        """Add title slide."""
        # Use blank layout and add shapes manually
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        
        # Format title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_run = title_para.runs[0]
        title_run.font.size = Pt(44)
        title_run.font.bold = True
        title_run.font.color.rgb = self.title_color

    def _add_content_slide(self, prs: Presentation, content: SlideContent) -> None:
        """Add content slide with title and bullets."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_left = Inches(0.5)
        title_top = Inches(0.5)
        title_width = Inches(9)
        title_height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.text = content.title
        
        # Format title
        title_para = title_frame.paragraphs[0]
        title_run = title_para.runs[0]
        title_run.font.size = Pt(32)
        title_run.font.bold = True
        title_run.font.color.rgb = self.title_color
        
        # Add content box
        content_left = Inches(0.5)
        content_top = Inches(1.5)
        content_width = Inches(9)
        content_height = Inches(5.5)
        
        content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        # Add bullet points
        for i, bullet in enumerate(content.content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = self.text_color
            p.space_before = Pt(12)
        
        # Add notes if provided
        if content.notes:
            notes_slide = slide.notes_slide
            notes_frame = notes_slide.notes_text_frame
            notes_frame.text = content.notes

    async def generate_pdf(
        self,
        title: str,
        slides: List[SlideContent],
    ) -> bytes:
        """
        Generate presentation in PDF format.
        
        Args:
            title: Presentation title
            slides: List of slide contents
            
        Returns:
            PDF file as bytes
        """
        # Generate HTML presentation
        html_content = self._generate_html_presentation(title, slides)
        
        # Convert to PDF using Playwright
        async with async_playwright() as p:
            browser = await p.chromium.launch()
            page = await browser.new_page()
            await page.set_content(html_content)
            
            # Set viewport for 16:9 aspect ratio
            await page.set_viewport_size({"width": 1280, "height": 720})
            
            pdf_bytes = await page.pdf(
                format='A4',
                landscape=True,
                margin={
                    'top': '0.5in',
                    'right': '0.5in',
                    'bottom': '0.5in',
                    'left': '0.5in',
                },
                print_background=True,
            )
            await browser.close()
        
        return pdf_bytes

    def _generate_html_presentation(self, title: str, slides: List[SlideContent]) -> str:
        """Generate HTML for PDF conversion."""
        html = """
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <style>
        @page {
            size: A4 landscape;
            margin: 0;
        }
        body {
            font-family: 'Calibri', 'Arial', sans-serif;
            margin: 0;
            padding: 0;
        }
        .slide {
            width: 100%;
            height: 100vh;
            padding: 40px 60px;
            box-sizing: border-box;
            page-break-after: always;
            display: flex;
            flex-direction: column;
        }
        .slide:last-child {
            page-break-after: auto;
        }
        .title-slide {
            justify-content: center;
            align-items: center;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
        }
        .title-slide h1 {
            font-size: 48pt;
            margin: 0;
            text-align: center;
        }
        .content-slide h2 {
            color: #1f4e79;
            font-size: 32pt;
            margin: 0 0 30px 0;
            border-bottom: 3px solid #1f4e79;
            padding-bottom: 10px;
        }
        .content-slide ul {
            list-style: none;
            padding: 0;
            margin: 0;
        }
        .content-slide li {
            font-size: 18pt;
            color: #333;
            margin-bottom: 20px;
            padding-left: 30px;
            position: relative;
        }
        .content-slide li:before {
            content: "●";
            color: #1f4e79;
            font-size: 14pt;
            position: absolute;
            left: 0;
        }
    </style>
</head>
<body>
"""
        
        # Title slide
        html += f"""
    <div class="slide title-slide">
        <h1>{title}</h1>
    </div>
"""
        
        # Content slides
        for slide_content in slides:
            html += f"""
    <div class="slide content-slide">
        <h2>{slide_content.title}</h2>
        <ul>
"""
            for bullet in slide_content.content:
                html += f"            <li>{bullet}</li>\n"
            
            html += """        </ul>
    </div>
"""
        
        html += """
</body>
</html>
"""
        return html

